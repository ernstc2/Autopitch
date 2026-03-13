"""Streamlit web UI for Autopitch.

Upload an Excel financial workbook and download a consulting-quality PPTX deck.
All pipeline logic lives in autopitch/pipeline.py — this file only renders widgets.
"""

import hashlib
import time
from io import BytesIO
from pathlib import Path

import openpyxl
from pptx import Presentation
import streamlit as st
from dotenv import load_dotenv

load_dotenv()

from autopitch.pipeline import run_pipeline  # noqa: E402
from autopitch.models import ValidationError  # noqa: E402

# ---------------------------------------------------------------------------
# Page config
# ---------------------------------------------------------------------------

st.set_page_config(
    page_title="Autopitch",
    page_icon="A",
    layout="wide",
)

# ---------------------------------------------------------------------------
# Custom CSS
# ---------------------------------------------------------------------------

st.markdown(
    """
    <style>
    /* ── Import fonts ─────────────────────────────────────── */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');
    @import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Mono:wght@400;500;600&display=swap');

    /* ── Animations ───────────────────────────────────────── */
    @keyframes fade-in {
        from { opacity: 0; transform: translateY(6px); }
        to { opacity: 1; transform: translateY(0); }
    }
    @media (prefers-reduced-motion: reduce) {
        *, *::before, *::after {
            animation-duration: 0.01ms !important;
            transition-duration: 0.01ms !important;
        }
    }

    /* ── Global overrides ─────────────────────────────────── */
    html, body, [class*="css"], .stMarkdown, .stMarkdown p {
        font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif !important;
    }

    /* Hide Streamlit chrome */
    h1 a, h2 a, h3 a { display: none !important; }
    header[data-testid="stHeader"] { background: transparent !important; }

    /* Layout */
    .block-container {
        max-width: 1060px !important;
        padding-top: 1.25rem !important;
        padding-bottom: 2rem !important;
    }

    /* ── Top bar (system info strip) ──────────────────────── */
    .top-bar {
        display: flex;
        align-items: center;
        justify-content: space-between;
        padding: 0.4rem 0;
        margin-bottom: 0.75rem;
        border-bottom: 1px solid rgba(0, 131, 138, 0.1);
        font-family: 'IBM Plex Mono', monospace;
        font-size: 0.65rem;
        letter-spacing: 0.05em;
        color: #4a6380;
        text-transform: uppercase;
        animation: fade-in 0.3s ease-out;
    }
    .top-bar-left {
        display: flex;
        align-items: center;
        gap: 1rem;
    }
    .top-bar-item {
        display: flex;
        align-items: center;
        gap: 0.3rem;
    }
    .top-bar-label {
        color: #3d4f63;
    }
    .top-bar-val {
        color: #7a8ba3;
    }
    .top-bar-dot {
        width: 5px;
        height: 5px;
        border-radius: 50%;
        background: #00838A;
        display: inline-block;
    }

    /* ── App header ───────────────────────────────────────── */
    .app-header {
        margin-bottom: 1.5rem;
        animation: fade-in 0.35s ease-out;
    }
    .app-name {
        font-size: 1.4rem !important;
        font-weight: 800 !important;
        color: #e8edf4 !important;
        letter-spacing: 0.06em !important;
        text-transform: uppercase !important;
        margin: 0 0 0.3rem 0 !important;
    }
    .app-name-accent {
        color: #00838A;
    }
    .app-desc {
        font-size: 0.85rem;
        color: #6b7f96;
        line-height: 1.5;
    }

    /* ── Section panels ───────────────────────────────────── */
    .panel {
        border: 1px solid rgba(0, 131, 138, 0.1);
        border-left: 3px solid rgba(0, 131, 138, 0.4);
        border-radius: 4px;
        margin-bottom: 1.25rem;
        overflow: hidden;
        animation: fade-in 0.4s ease-out;
        transition: border-color 0.2s ease;
    }
    .panel:hover {
        border-color: rgba(0, 131, 138, 0.2);
        border-left-color: rgba(0, 131, 138, 0.6);
    }
    .panel-header {
        display: flex;
        align-items: center;
        justify-content: space-between;
        padding: 0.45rem 0.75rem;
        background: rgba(0, 131, 138, 0.07);
        border-bottom: 1px solid rgba(0, 131, 138, 0.1);
    }
    .panel-title {
        font-family: 'IBM Plex Mono', monospace;
        font-size: 0.68rem;
        font-weight: 600;
        color: #00838A;
        text-transform: uppercase;
        letter-spacing: 0.08em;
    }
    .panel-tag {
        font-family: 'IBM Plex Mono', monospace;
        font-size: 0.6rem;
        color: #3d4f63;
        letter-spacing: 0.04em;
    }
    .panel-body {
        padding: 0.75rem;
    }

    /* ── Section text ─────────────────────────────────────── */
    .section-heading {
        font-size: 1rem !important;
        font-weight: 700 !important;
        color: #e8edf4 !important;
        margin-bottom: 0.15rem !important;
    }
    .section-sub {
        font-size: 0.8rem;
        color: #6b7f96;
        margin-top: 0.1rem;
        margin-bottom: 0.5rem;
        line-height: 1.5;
    }

    /* ── Stat readout (data-field grid) ────────────────────── */
    .stat-readout {
        display: flex;
        align-items: flex-start;
        gap: 1.5rem;
        padding: 0.5rem 0.6rem;
        background: rgba(212, 146, 11, 0.04);
        border: 1px solid rgba(212, 146, 11, 0.12);
        border-left: 3px solid rgba(212, 146, 11, 0.35);
        border-radius: 4px;
        margin-top: 0.25rem;
    }
    .stat-field {
        display: flex;
        flex-direction: column;
        gap: 0.15rem;
    }
    .stat-label {
        font-family: 'IBM Plex Mono', monospace;
        font-size: 0.58rem;
        color: #4a6380;
        text-transform: uppercase;
        letter-spacing: 0.06em;
    }
    .stat-value {
        font-family: 'IBM Plex Mono', monospace;
        font-size: 0.85rem;
        font-weight: 600;
        color: #E0A81E;
        letter-spacing: 0.02em;
    }
    .stat-value.teal {
        color: #00838A;
    }

    /* .stat-pill — kept for test compatibility */
    .stat-pill {
        display: inline-block;
        background: rgba(212, 146, 11, 0.1);
        border: 1px solid rgba(212, 146, 11, 0.2);
        border-radius: 4px;
        padding: 0.25rem 0.6rem;
        font-family: 'IBM Plex Mono', monospace;
        font-size: 0.75rem;
        color: #E0A81E;
        margin-top: 0.4rem;
    }

    /* ── Deck preview ─────────────────────────────────────── */
    .deck-preview {
        margin-top: 0.4rem;
    }
    .deck-preview-row {
        display: flex;
        align-items: baseline;
        padding: 0.3rem 0.5rem;
        border-bottom: 1px solid rgba(0, 131, 138, 0.05);
        transition: background 0.15s ease;
    }
    .deck-preview-row:hover {
        background: rgba(0, 131, 138, 0.03);
    }
    .deck-preview-row:last-child {
        border-bottom: none;
    }
    .deck-preview-num {
        font-family: 'IBM Plex Mono', monospace;
        color: rgba(0, 131, 138, 0.6);
        font-size: 0.68rem;
        font-weight: 500;
        min-width: 2rem;
        font-variant-numeric: tabular-nums;
    }
    .deck-preview-title {
        color: #a0b4cc;
        font-size: 0.8rem;
        line-height: 1.4;
    }

    /* ── Streamlit widget overrides ────────────────────────── */
    .stButton > button[kind="primary"],
    .stButton > button[data-testid="stBaseButton-primary"] {
        font-weight: 600 !important;
        font-size: 0.8rem !important;
        letter-spacing: 0.04em !important;
        text-transform: uppercase !important;
        border-radius: 4px !important;
        padding: 0.45rem 1.25rem !important;
        transition: box-shadow 0.2s ease, transform 0.2s ease !important;
    }
    .stButton > button[kind="primary"]:hover,
    .stButton > button[data-testid="stBaseButton-primary"]:hover {
        box-shadow: 0 2px 12px rgba(0, 131, 138, 0.2) !important;
        transform: translateY(-1px);
    }

    .stDownloadButton > button {
        font-weight: 600 !important;
        font-size: 0.76rem !important;
        letter-spacing: 0.03em !important;
        text-transform: uppercase !important;
        border-radius: 4px !important;
        border-left: 3px solid rgba(212, 146, 11, 0.4) !important;
        transition: box-shadow 0.2s ease, border-color 0.2s ease !important;
    }
    .stDownloadButton > button:hover {
        border-left-color: rgba(212, 146, 11, 0.7) !important;
    }

    [data-testid="stFileUploader"] section {
        border: 1px dashed rgba(0, 131, 138, 0.18) !important;
        border-radius: 4px !important;
        background: rgba(0, 131, 138, 0.02) !important;
        transition: border-color 0.2s ease, background 0.2s ease !important;
    }
    [data-testid="stFileUploader"] section:hover {
        border-color: rgba(0, 131, 138, 0.3) !important;
        background: rgba(0, 131, 138, 0.04) !important;
    }

    .streamlit-expanderHeader {
        font-size: 0.8rem !important;
        font-weight: 500 !important;
        color: #7a8ba3 !important;
    }

    /* ── Divider ──────────────────────────────────────────── */
    hr {
        border-color: rgba(0, 131, 138, 0.08) !important;
        margin: 0.5rem 0 !important;
    }

    /* ── Footer ───────────────────────────────────────────── */
    .footer {
        font-family: 'IBM Plex Mono', monospace;
        color: #3d4f63;
        font-size: 0.68rem;
        padding: 0.75rem 0 0.5rem 0;
        border-top: 1px solid rgba(0, 131, 138, 0.08);
        margin-top: 1.5rem;
        line-height: 1.6;
        letter-spacing: 0.02em;
    }
    .footer-sep {
        color: #1e3350;
        margin: 0 0.4rem;
    }
    .footer-tech {
        color: rgba(0, 131, 138, 0.7);
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------------
# Session state initialization — must come before any widget rendering
# ---------------------------------------------------------------------------

for _key in ("demo_pptx", "demo_elapsed", "demo_slides", "demo_titles",
             "upload_pptx", "upload_hash", "upload_titles"):
    if _key not in st.session_state:
        st.session_state[_key] = None

# ---------------------------------------------------------------------------
# Cached demo pipeline — first click calls Claude; subsequent serve cache
# ---------------------------------------------------------------------------


@st.cache_data(show_spinner=False)
def _run_demo_pipeline() -> bytes:
    """Cached: returns PPTX bytes for the bundled Apple demo file.

    @st.cache_data ensures repeated 'Try the Demo' clicks do not re-call the
    Anthropic API — only the first invocation hits the network.
    """
    return run_pipeline(Path("demo/apple_financials.xlsx"))


# ---------------------------------------------------------------------------
# Excel template builder — used in upload section below
# ---------------------------------------------------------------------------


def _build_template_xlsx() -> bytes:
    """Build a minimal openpyxl workbook with the required sheet/column structure.

    Returns:
        Valid XLSX bytes (ZIP magic PK at offset 0).
    """
    wb = openpyxl.Workbook()

    # P&L sheet
    ws_pl = wb.active
    ws_pl.title = "P&L"
    ws_pl.append(["", "FY2022", "FY2023", "FY2024"])
    for label in [
        "Revenue",
        "COGS",
        "Operating Income",
        "Depreciation & Amortization",
        "Net Income",
    ]:
        ws_pl.append([label, "", "", ""])

    # Balance Sheet
    ws_bs = wb.create_sheet("Balance Sheet")
    ws_bs.append(["", "FY2022", "FY2023", "FY2024"])
    for label in [
        "Current Assets",
        "Current Liabilities",
        "Total Assets",
        "Total Debt",
        "Total Shareholders Equity",
    ]:
        ws_bs.append([label, "", "", ""])

    # Cash Flow
    ws_cf = wb.create_sheet("Cash Flow")
    ws_cf.append(["", "FY2022", "FY2023", "FY2024"])
    for label in ["Operating Cash Flow", "Capital Expenditures"]:
        ws_cf.append([label, "", "", ""])

    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Deck preview helpers
# ---------------------------------------------------------------------------


def _extract_slide_titles(pptx_bytes: bytes) -> list[str]:
    """Parse PPTX bytes and extract the title text from each slide's header.

    Looks for the first text-bearing shape on each slide (the header title
    added by _add_header in deck.py). Falls back to 'Slide N' if no text found.
    """
    prs = Presentation(BytesIO(pptx_bytes))
    titles = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip()
                break
        titles.append(title or f"Slide {i}")
    return titles


def _render_deck_preview(titles: list[str]) -> None:
    """Render a compact slide table of contents using custom HTML."""
    with st.expander(f"Deck contents \u2014 {len(titles)} slides", expanded=False):
        rows_html = ""
        for i, title in enumerate(titles, start=1):
            rows_html += (
                f'<div class="deck-preview-row">'
                f'<span class="deck-preview-num">{i}.</span>'
                f'<span class="deck-preview-title">{title}</span>'
                f'</div>'
            )
        st.markdown(
            f'<div class="deck-preview">{rows_html}</div>',
            unsafe_allow_html=True,
        )


# ═══════════════════════════════════════════════════════════════════════════
# APP HEADER — compact, left-aligned, no marketing
# ═══════════════════════════════════════════════════════════════════════════

st.markdown(
    '<div class="top-bar">'
    '<div class="top-bar-left">'
    '<div class="top-bar-item">'
    '<span class="top-bar-dot"></span>'
    '<span class="top-bar-val">Autopitch v1.0</span>'
    '</div>'
    '<div class="top-bar-item">'
    '<span class="top-bar-label">Engine</span>'
    '<span class="top-bar-val">Claude</span>'
    '</div>'
    '<div class="top-bar-item">'
    '<span class="top-bar-label">Output</span>'
    '<span class="top-bar-val">PPTX</span>'
    '</div>'
    '</div>'
    '<div class="top-bar-item">'
    '<span class="top-bar-dot"></span>'
    '<span style="color:#00838A">Ready</span>'
    '</div>'
    '</div>',
    unsafe_allow_html=True,
)

st.markdown(
    '<div class="app-header">'
    '<div class="app-name">'
    'Auto<span class="app-name-accent">pitch</span>'
    '</div>'
    '<div class="app-desc">'
    "Excel financials in, consulting-quality pitch deck out."
    "</div>"
    "</div>",
    unsafe_allow_html=True,
)

# ═══════════════════════════════════════════════════════════════════════════
# DEMO SECTION
# ═══════════════════════════════════════════════════════════════════════════

st.markdown(
    '<div class="panel"><div class="panel-header">'
    '<span class="panel-title">Demo</span>'
    '<span class="panel-tag">AAPL FY2020-2024</span>'
    '</div><div class="panel-body">'
    '<div class="section-heading">Try the Demo</div>'
    '<div class="section-sub">'
    "See a sample deck built from Apple\u2019s public FY2020\u20132024 financials."
    "</div>"
    "</div></div>",
    unsafe_allow_html=True,
)

if st.button("Generate Apple Demo Deck", type="primary"):
    if st.session_state["demo_pptx"] is None:
        t0 = time.perf_counter()
        with st.spinner("Analyzing financials and assembling deck\u2026"):
            pptx_bytes = _run_demo_pipeline()
        st.session_state["demo_pptx"] = pptx_bytes
        st.session_state["demo_elapsed"] = time.perf_counter() - t0
        st.session_state["demo_titles"] = _extract_slide_titles(pptx_bytes)
        st.session_state["demo_slides"] = len(st.session_state["demo_titles"])

# Render download button and stats based on session_state (persists across reruns)
if st.session_state["demo_pptx"] is not None:
    col1, col2 = st.columns([2, 3])
    with col1:
        st.download_button(
            label="Download Demo Deck",
            data=st.session_state["demo_pptx"],
            file_name="autopitch_apple_demo.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            on_click="ignore",
        )
    with col2:
        elapsed = st.session_state["demo_elapsed"]
        slides = st.session_state["demo_slides"]
        st.markdown(
            f'<div class="stat-readout">'
            f'<div class="stat-field">'
            f'<span class="stat-label">Slides</span>'
            f'<span class="stat-value">{slides}</span>'
            f'</div>'
            f'<div class="stat-field">'
            f'<span class="stat-label">Build Time</span>'
            f'<span class="stat-value teal">{elapsed:.1f}s</span>'
            f'</div>'
            f'<div class="stat-field">'
            f'<span class="stat-label">Status</span>'
            f'<span class="stat-value teal">READY</span>'
            f'</div>'
            f'</div>'
            f'<span class="stat-pill" style="display:none">{slides} slides &middot; {elapsed:.1f}s</span>',
            unsafe_allow_html=True,
        )
    _render_deck_preview(st.session_state["demo_titles"])

st.divider()

# ═══════════════════════════════════════════════════════════════════════════
# UPLOAD SECTION
# ═══════════════════════════════════════════════════════════════════════════

st.markdown(
    '<div class="panel"><div class="panel-header">'
    '<span class="panel-title">Upload</span>'
    '<span class="panel-tag">Custom Financials</span>'
    '</div><div class="panel-body">'
    '<div class="section-heading">Upload Your Own Data</div>'
    '<div class="section-sub">'
    "Three sheets required: P&L, Balance Sheet, and Cash Flow."
    "</div>"
    "</div></div>",
    unsafe_allow_html=True,
)

uploaded = st.file_uploader(
    "Excel workbook (.xlsx)",
    type=["xlsx"],
    help="Not sure about the format? Expand the format guide or download the template below.",
)

col_help1, col_help2 = st.columns(2)
with col_help1:
    with st.expander("Format guide \u2014 required sheets and rows"):
        st.markdown(
            """
**Three sheets required (exact names):** `P&L`, `Balance Sheet`, `Cash Flow`

**Column layout:** Column A = row labels, columns B+ = fiscal years (`FY2022`, `FY2023`, ...)

**P&L required rows:** Revenue, COGS, Operating Income, Depreciation & Amortization, Net Income

**Balance Sheet required rows:** Current Assets, Current Liabilities, Total Assets,
Total Debt, Total Shareholders Equity

**Cash Flow required rows:** Operating Cash Flow, Capital Expenditures

All values must be positive numbers. Use consistent units across all sheets (e.g., millions: enter 394,328 for $394.3B).
            """
        )
with col_help2:
    st.download_button(
        label="Download Excel Template",
        data=_build_template_xlsx(),
        file_name="autopitch_template.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        on_click="ignore",
    )

if uploaded is not None:
    if st.button("Generate Deck", type="primary"):
        uploaded.seek(0)
        file_hash = hashlib.md5(uploaded.read()).hexdigest()
        uploaded.seek(0)
        if file_hash == st.session_state["upload_hash"] and st.session_state["upload_pptx"] is not None:
            st.success("This file was already processed \u2014 ready to download.")
        else:
            with st.spinner("Analyzing financials and assembling deck\u2026"):
                try:
                    pptx_bytes = run_pipeline(uploaded)
                    st.session_state["upload_pptx"] = pptx_bytes
                    st.session_state["upload_hash"] = file_hash
                    st.session_state["upload_titles"] = _extract_slide_titles(pptx_bytes)
                    slides = len(st.session_state["upload_titles"])
                    st.success(f"Deck ready \u2014 {slides} slides generated.")
                except ValidationError as e:
                    if len(e.errors) == 1:
                        st.error(e.errors[0])
                    else:
                        bullets = "\n".join(f"- {err}" for err in e.errors)
                        st.error(f"Found {len(e.errors)} issues:\n\n{bullets}")

# Render upload download button and preview based on session_state (persists across reruns)
if st.session_state["upload_pptx"] is not None:
    output_name = (uploaded.name.replace(".xlsx", "") + ".pptx") if uploaded else "autopitch_deck.pptx"
    st.download_button(
        label="Download Deck",
        data=st.session_state["upload_pptx"],
        file_name=output_name,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        on_click="ignore",
    )
    if st.session_state["upload_titles"] is not None:
        _render_deck_preview(st.session_state["upload_titles"])

# ═══════════════════════════════════════════════════════════════════════════
# FOOTER
# ═══════════════════════════════════════════════════════════════════════════

st.markdown(
    '<div class="footer">'
    "<span class='footer-tech'>Python</span>"
    "<span class='footer-sep'>/</span>"
    "<span class='footer-tech'>Claude API</span>"
    "<span class='footer-sep'>/</span>"
    "<span class='footer-tech'>python-pptx</span>"
    "<span class='footer-sep'>/</span>"
    "<span class='footer-tech'>matplotlib</span>"
    "<span class='footer-sep'>&mdash;</span>Built by Chris"
    "</div>",
    unsafe_allow_html=True,
)
