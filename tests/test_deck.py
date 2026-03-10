"""
Failing test stubs for DECK-01..08.
All tests are skipped until autopitch.deck is implemented.
"""
import pytest
from pptx.dml.color import RGBColor

try:
    from autopitch.deck import build_deck
    _DECK_AVAILABLE = True
except ImportError:
    _DECK_AVAILABLE = False

pytestmark = pytest.mark.skipif(not _DECK_AVAILABLE, reason="autopitch.deck not yet implemented")


def _get_all_text(slide):
    """Collect all textbox strings from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return " ".join(texts)


def _get_header_text(slide):
    """Return the text of the first textbox in the slide (assumed to be the header)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape.text_frame.text
    return ""


def test_title_slide(minimal_valid_data, metrics_from_minimal):
    """DECK-01: First slide contains the company name."""
    prs = build_deck(minimal_valid_data, metrics_from_minimal)
    slide = prs.slides[0]
    all_text = _get_all_text(slide)
    assert minimal_valid_data.company_name in all_text, (
        f"Company name '{minimal_valid_data.company_name}' not found on title slide. "
        f"Found text: '{all_text}'"
    )


def test_exec_summary_slide(minimal_valid_data, metrics_from_minimal):
    """DECK-02: Second slide contains non-empty text."""
    prs = build_deck(minimal_valid_data, metrics_from_minimal)
    assert len(prs.slides) >= 2, "Deck must have at least 2 slides"
    slide = prs.slides[1]
    all_text = _get_all_text(slide).strip()
    assert len(all_text) > 0, "Executive summary slide must contain non-empty text"


def test_pl_section_slides(minimal_valid_data, metrics_from_minimal):
    """DECK-03: Between 3 and 4 slides contain 'P&L' in their header text."""
    prs = build_deck(minimal_valid_data, metrics_from_minimal)
    pl_slides = [
        slide for slide in prs.slides
        if "P&L" in _get_all_text(slide)
    ]
    assert 3 <= len(pl_slides) <= 4, (
        f"Expected 3-4 P&L slides, found {len(pl_slides)}"
    )


def test_bs_section_slides(minimal_valid_data, metrics_from_minimal):
    """DECK-04: Between 2 and 3 slides contain 'Balance Sheet' in their header text."""
    prs = build_deck(minimal_valid_data, metrics_from_minimal)
    bs_slides = [
        slide for slide in prs.slides
        if "Balance Sheet" in _get_all_text(slide)
    ]
    assert 2 <= len(bs_slides) <= 3, (
        f"Expected 2-3 Balance Sheet slides, found {len(bs_slides)}"
    )


def test_cf_section_slides(minimal_valid_data, metrics_from_minimal):
    """DECK-05: Between 2 and 3 slides contain 'Cash Flow' in their header text."""
    prs = build_deck(minimal_valid_data, metrics_from_minimal)
    cf_slides = [
        slide for slide in prs.slides
        if "Cash Flow" in _get_all_text(slide)
    ]
    assert 2 <= len(cf_slides) <= 3, (
        f"Expected 2-3 Cash Flow slides, found {len(cf_slides)}"
    )


def test_kpi_slide(minimal_valid_data, metrics_from_minimal):
    """DECK-06: Exactly 1 slide contains 'KPI' or 'Scorecard' in its text."""
    prs = build_deck(minimal_valid_data, metrics_from_minimal)
    kpi_slides = [
        slide for slide in prs.slides
        if "KPI" in _get_all_text(slide) or "Scorecard" in _get_all_text(slide)
    ]
    assert len(kpi_slides) == 1, (
        f"Expected exactly 1 KPI/Scorecard slide, found {len(kpi_slides)}"
    )


def test_footer_on_all_slides(minimal_valid_data, metrics_from_minimal):
    """DECK-07: Every slide has at least one textbox containing the company name."""
    prs = build_deck(minimal_valid_data, metrics_from_minimal)
    company = minimal_valid_data.company_name
    for i, slide in enumerate(prs.slides):
        all_text = _get_all_text(slide)
        assert company in all_text, (
            f"Slide {i} missing company name '{company}' in footer. "
            f"Found text: '{all_text}'"
        )


def test_navy_header_color(minimal_valid_data, metrics_from_minimal):
    """DECK-08: Slide 0 has a shape with fill color == NAVY (#002B49)."""
    prs = build_deck(minimal_valid_data, metrics_from_minimal)
    slide = prs.slides[0]
    navy = RGBColor(0x00, 0x2B, 0x49)
    navy_found = False
    for shape in slide.shapes:
        try:
            fill = shape.fill
            if fill.type is not None:
                rgb = fill.fore_color.rgb
                if rgb == navy:
                    navy_found = True
                    break
        except Exception:
            continue
    assert navy_found, (
        f"No shape with navy fill (#002B49) found on slide 0"
    )
