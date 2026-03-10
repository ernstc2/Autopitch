"""
Deck assembly module for Autopitch.

build_deck(data, metrics) -> Presentation

Constructs a complete 11-slide PPTX with Big 4 aesthetics from FinancialData
and MetricsOutput.  All charts are rendered as matplotlib PNG images (via
charts.py) and embedded in the content area of each slide.
"""
from __future__ import annotations

import io
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from autopitch.models import FinancialData, MetricsOutput
from autopitch.theme import (
    NAVY, TEAL, WHITE, DARK_GRAY, LIGHT_GRAY,
    FONT_HEADING, FONT_BODY,
    SLIDE_W, SLIDE_H,
    HEADER_HEIGHT,
    CONTENT_TOP, CONTENT_LEFT, CONTENT_W, CONTENT_H,
    FOOTER_TOP, FOOTER_HEIGHT,
    SZ_SLIDE_TITLE, SZ_SECTION_TITLE, SZ_BODY, SZ_FOOTER,
)
from autopitch.charts import bar_chart, line_chart, waterfall_chart, kpi_scorecard


# ---------------------------------------------------------------------------
# Private slide helpers
# ---------------------------------------------------------------------------

def _add_slide(prs: Presentation):
    """Add a blank slide using layout index 6 (Blank layout)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _add_header(slide, title: str) -> None:
    """Add a navy header bar spanning the full width, with white title text."""
    shapes = slide.shapes

    # Navy color block (rectangle shape type = 1)
    bar = shapes.add_shape(1, 0, 0, SLIDE_W, HEADER_HEIGHT)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()  # Remove default border

    # Title text over the bar
    txBox = shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12.0), Inches(0.75))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEADING
    run.font.size = SZ_SLIDE_TITLE
    run.font.bold = True
    run.font.color.rgb = WHITE


def _add_footer(slide, company: str, slide_num: int, total: int) -> None:
    """Add a navy footer bar with company name (left) and slide number (right)."""
    shapes = slide.shapes

    # Navy footer bar
    bar = shapes.add_shape(1, 0, FOOTER_TOP, SLIDE_W, FOOTER_HEIGHT)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Company name + date on the left
    txL = shapes.add_textbox(
        Inches(0.4), FOOTER_TOP + Inches(0.06), Inches(9.0), FOOTER_HEIGHT
    )
    tf_l = txL.text_frame
    pL = tf_l.paragraphs[0]
    runL = pL.add_run()
    runL.text = f"{company}  |  {date.today().strftime('%B %Y')}"
    runL.font.name = FONT_BODY
    runL.font.size = SZ_FOOTER
    runL.font.color.rgb = WHITE

    # Slide number on the right
    txR = shapes.add_textbox(
        Inches(11.0), FOOTER_TOP + Inches(0.06), Inches(2.0), FOOTER_HEIGHT
    )
    tf_r = txR.text_frame
    pR = tf_r.paragraphs[0]
    pR.alignment = PP_ALIGN.RIGHT
    runR = pR.add_run()
    runR.text = f"{slide_num} / {total}"
    runR.font.name = FONT_BODY
    runR.font.size = SZ_FOOTER
    runR.font.color.rgb = WHITE


def _embed_chart(slide, buf: io.BytesIO) -> None:
    """Embed a matplotlib PNG BytesIO into the slide content area."""
    buf.seek(0)  # Guard: rewind in case cursor moved
    slide.shapes.add_picture(buf, CONTENT_LEFT, CONTENT_TOP, CONTENT_W, CONTENT_H)


def _add_text_block(
    slide,
    text: str,
    top=CONTENT_TOP,
    left=CONTENT_LEFT,
    width=CONTENT_W,
    height=CONTENT_H,
    size=SZ_BODY,
) -> None:
    """Add a textbox with word wrap enabled."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = size
    run.font.color.rgb = DARK_GRAY


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def build_deck(data: FinancialData, metrics: MetricsOutput) -> Presentation:
    """Assemble a complete 11-slide PPTX from FinancialData and MetricsOutput.

    Assembly sequence:
      1. Create Presentation; set 16:9 dimensions
      2. Generate ALL charts upfront (fail fast before PPTX mutation)
      3. Build each slide in section order
      4. Footer pass: iterate all slides with known total count
      5. Return Presentation ready for .save()
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    company = data.company_name
    years = metrics.years

    # ------------------------------------------------------------------
    # Step 1: Extract data arrays (needed by charts and executive summary)
    # ------------------------------------------------------------------

    def _get_values(rows: dict, key: str, yr_list: list) -> list[float]:
        """Return float values for a row key across yr_list; default 0 for None."""
        row = rows.get(key, {})
        return [float(row.get(yr) or 0) for yr in yr_list]

    # P&L data
    revenue_vals = _get_values(data.pl.rows, "Revenue", years)
    cogs_vals = _get_values(data.pl.rows, "COGS", years)
    op_income_vals = _get_values(data.pl.rows, "Operating Income", years)
    da_vals = _get_values(data.pl.rows, "Depreciation & Amortization", years)
    net_income_vals = _get_values(data.pl.rows, "Net Income", years)

    # Balance sheet data
    total_assets_vals = _get_values(data.balance_sheet.rows, "Total Assets", years)
    total_debt_vals = _get_values(data.balance_sheet.rows, "Total Debt", years)

    # Cash flow data
    ocf_vals = _get_values(data.cash_flow.rows, "Operating Cash Flow", years)
    capex_vals = _get_values(data.cash_flow.rows, "Capital Expenditures", years)
    fcf_vals = [
        float(metrics.free_cash_flow.get(yr) or 0) for yr in years
    ]

    # Margin series (percent values)
    gross_margin_vals = [float(metrics.gross_margin.get(yr) or 0) for yr in years]
    ebitda_margin_vals = [float(metrics.ebitda_margin.get(yr) or 0) for yr in years]
    net_margin_vals = [float(metrics.net_margin.get(yr) or 0) for yr in years]

    # Working capital / liquidity series
    wc_vals = [float(metrics.working_capital.get(yr) or 0) for yr in years]
    cr_vals = [float(metrics.current_ratio.get(yr) or 0) for yr in years]

    # ------------------------------------------------------------------
    # Step 2: Build waterfall data for the most recent year P&L bridge
    # Most recent year bridge: Revenue → -COGS → Gross Profit → -OpEx → EBITDA → -Taxes/Other → Net Income
    # ------------------------------------------------------------------
    yr_last = years[-1]
    rev_lyr = float(data.pl.rows.get("Revenue", {}).get(yr_last) or 0)
    cogs_lyr = float(data.pl.rows.get("COGS", {}).get(yr_last) or 0)
    op_income_lyr = float(data.pl.rows.get("Operating Income", {}).get(yr_last) or 0)
    da_lyr = float(data.pl.rows.get("Depreciation & Amortization", {}).get(yr_last) or 0)
    net_income_lyr = float(data.pl.rows.get("Net Income", {}).get(yr_last) or 0)

    gross_profit_lyr = rev_lyr - cogs_lyr
    ebitda_lyr = op_income_lyr + da_lyr
    # Other = EBITDA - Net Income (taxes + interest net)
    other_lyr = ebitda_lyr - net_income_lyr

    wf_labels = ["Revenue", "- COGS", "Gross\nProfit", "- OpEx", "EBITDA", "- Other", "Net\nIncome"]
    wf_values = [
        rev_lyr,
        -cogs_lyr,
        gross_profit_lyr,
        -(ebitda_lyr - gross_profit_lyr),  # OpEx = Gross Profit - EBITDA (negative delta)
        ebitda_lyr,
        -other_lyr,
        net_income_lyr,
    ]

    # ------------------------------------------------------------------
    # Step 3: Build KPI scorecard metrics dict (most recent year)
    # ------------------------------------------------------------------
    yr = years[-1]
    kpi_data = {
        "Revenue Growth":  f"{metrics.revenue_growth.get(yr, 0) or 0:.1f}%",
        "Gross Margin":    f"{metrics.gross_margin.get(yr, 0) or 0:.1f}%",
        "EBITDA Margin":   f"{metrics.ebitda_margin.get(yr, 0) or 0:.1f}%",
        "Net Margin":      f"{metrics.net_margin.get(yr, 0) or 0:.1f}%",
        "Current Ratio":   f"{metrics.current_ratio.get(yr, 0) or 0:.2f}x",
        "D/E Ratio":       (
            f"{metrics.debt_to_equity.get(yr, 0) or 0:.2f}x"
            if metrics.debt_to_equity.get(yr) is not None
            else "N/A"
        ),
        "ROE":             (
            f"{metrics.roe.get(yr, 0) or 0:.1f}%"
            if metrics.roe.get(yr) is not None
            else "N/A"
        ),
        "Free Cash Flow":  f"${(metrics.free_cash_flow.get(yr, 0) or 0) / 1000:.0f}B",
    }

    # ------------------------------------------------------------------
    # Step 4: Executive summary text (pre-computed metric values)
    # ------------------------------------------------------------------
    most_recent = years[-1]
    rev_mr = data.pl.rows.get("Revenue", {}).get(most_recent)
    net_m_mr = metrics.net_margin.get(most_recent)
    fcf_mr = metrics.free_cash_flow.get(most_recent)
    rev_growth_mr = metrics.revenue_growth.get(most_recent)
    gross_m_mr = metrics.gross_margin.get(most_recent)

    def _fmt_val(v, fmt=".0f", unit="$", divisor=1000, suffix="B") -> str:
        if v is None:
            return "N/A"
        return f"{unit}{v / divisor:{fmt}}{suffix}"

    def _fmt_pct(v) -> str:
        return "N/A" if v is None else f"{v:.1f}%"

    period_str = f"{years[0]} \u2013 {years[-1]}" if len(years) > 1 else years[0]
    exec_bullets = "\n".join([
        f"Revenue ({most_recent}): {_fmt_val(rev_mr)}",
        f"Net Margin ({most_recent}): {_fmt_pct(net_m_mr)}",
        f"Gross Margin ({most_recent}): {_fmt_pct(gross_m_mr)}",
        f"Revenue Growth ({most_recent}): {_fmt_pct(rev_growth_mr)}",
        f"Free Cash Flow ({most_recent}): {_fmt_val(fcf_mr)}",
    ])

    # ------------------------------------------------------------------
    # Step 5: Generate ALL charts upfront (fail fast before PPTX mutation)
    # ------------------------------------------------------------------
    charts: dict[str, io.BytesIO] = {}

    # Slide 3: P&L Revenue Trends
    charts["pl_revenue"] = bar_chart(
        years, revenue_vals,
        title=f"Revenue Trends ({period_str})",
        ylabel="Revenue (USD)"
    )

    # Slide 4: P&L Margin Analysis
    charts["pl_margins"] = line_chart(
        years,
        {
            "Gross Margin": gross_margin_vals,
            "EBITDA Margin": ebitda_margin_vals,
            "Net Margin": net_margin_vals,
        },
        title="Margin Analysis",
        ylabel="Margin (%)"
    )

    # Slide 5: P&L Earnings Bridge (waterfall)
    charts["pl_bridge"] = waterfall_chart(
        wf_labels, wf_values,
        title=f"P&L Bridge ({yr_last})"
    )

    # Slide 6: Balance Sheet Asset Composition
    charts["bs_assets"] = bar_chart(
        years, total_assets_vals,
        title="Total Assets",
        ylabel="Assets (USD)"
    )

    # Slide 7: Balance Sheet Working Capital & Liquidity
    charts["bs_wc"] = line_chart(
        years,
        {
            "Working Capital": wc_vals,
        },
        title="Working Capital Trend",
        ylabel="Working Capital (USD)"
    )

    # Slide 8: Balance Sheet Leverage
    charts["bs_leverage"] = bar_chart(
        years, total_debt_vals,
        title="Total Debt",
        ylabel="Debt (USD)"
    )

    # Slide 9: Cash Flow OCF vs FCF
    charts["cf_fcf"] = bar_chart(
        years, fcf_vals,
        title="Free Cash Flow",
        ylabel="FCF (USD)"
    )

    # Slide 10: Cash Flow Trend
    charts["cf_trend"] = line_chart(
        years,
        {
            "Free Cash Flow": fcf_vals,
        },
        title="Cash Flow Trend",
        ylabel="USD"
    )

    # Slide 11: KPI Scorecard
    charts["kpi"] = kpi_scorecard(
        kpi_data,
        title="KPI Scorecard"
    )

    # ------------------------------------------------------------------
    # Step 6: Slide 1 — Title slide (full-navy background, no standard header)
    # ------------------------------------------------------------------
    slide_title = _add_slide(prs)

    # Full-slide navy background rectangle
    bg = slide_title.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Company name (large, centered horizontally)
    tx_company = slide_title.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(11.33), Inches(1.2)
    )
    tf_c = tx_company.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    run_c = p_c.add_run()
    run_c.text = company
    run_c.font.name = FONT_HEADING
    run_c.font.size = Pt(36)
    run_c.font.bold = True
    run_c.font.color.rgb = WHITE

    # Subtitle
    tx_sub = slide_title.shapes.add_textbox(
        Inches(1.0), Inches(3.3), Inches(11.33), Inches(0.6)
    )
    tf_sub = tx_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    run_sub = p_sub.add_run()
    run_sub.text = "Financial Performance Review"
    run_sub.font.name = FONT_HEADING
    run_sub.font.size = Pt(20)
    run_sub.font.color.rgb = WHITE

    # Period line
    tx_period = slide_title.shapes.add_textbox(
        Inches(1.0), Inches(3.95), Inches(11.33), Inches(0.45)
    )
    tf_period = tx_period.text_frame
    p_period = tf_period.paragraphs[0]
    p_period.alignment = PP_ALIGN.CENTER
    run_period = p_period.add_run()
    run_period.text = period_str
    run_period.font.name = FONT_BODY
    run_period.font.size = Pt(14)
    run_period.font.color.rgb = WHITE

    # Generation date
    tx_date = slide_title.shapes.add_textbox(
        Inches(1.0), Inches(4.45), Inches(11.33), Inches(0.45)
    )
    tf_date = tx_date.text_frame
    p_date = tf_date.paragraphs[0]
    p_date.alignment = PP_ALIGN.CENTER
    run_date = p_date.add_run()
    run_date.text = date.today().strftime("%B %d, %Y")
    run_date.font.name = FONT_BODY
    run_date.font.size = Pt(12)
    run_date.font.color.rgb = WHITE

    # ------------------------------------------------------------------
    # Slide 2 — Executive Summary
    # ------------------------------------------------------------------
    slide_exec = _add_slide(prs)
    _add_header(slide_exec, "Executive Summary")
    _add_text_block(slide_exec, exec_bullets)

    # ------------------------------------------------------------------
    # Slide 3 — P&L | Revenue Trends
    # ------------------------------------------------------------------
    slide_pl_rev = _add_slide(prs)
    _add_header(slide_pl_rev, "P&L | Revenue Trends")
    _embed_chart(slide_pl_rev, charts["pl_revenue"])

    # ------------------------------------------------------------------
    # Slide 4 — P&L | Margin Analysis
    # ------------------------------------------------------------------
    slide_pl_margin = _add_slide(prs)
    _add_header(slide_pl_margin, "P&L | Margin Analysis")
    _embed_chart(slide_pl_margin, charts["pl_margins"])

    # ------------------------------------------------------------------
    # Slide 5 — P&L | Earnings Bridge
    # ------------------------------------------------------------------
    slide_pl_bridge = _add_slide(prs)
    _add_header(slide_pl_bridge, "P&L | Earnings Bridge")
    _embed_chart(slide_pl_bridge, charts["pl_bridge"])

    # ------------------------------------------------------------------
    # Slide 6 — Balance Sheet | Asset Composition
    # ------------------------------------------------------------------
    slide_bs_assets = _add_slide(prs)
    _add_header(slide_bs_assets, "Balance Sheet | Asset Composition")
    _embed_chart(slide_bs_assets, charts["bs_assets"])

    # ------------------------------------------------------------------
    # Slide 7 — Balance Sheet | Working Capital
    # ------------------------------------------------------------------
    slide_bs_wc = _add_slide(prs)
    _add_header(slide_bs_wc, "Balance Sheet | Working Capital")
    _embed_chart(slide_bs_wc, charts["bs_wc"])

    # ------------------------------------------------------------------
    # Slide 8 — Balance Sheet | Leverage
    # ------------------------------------------------------------------
    slide_bs_lev = _add_slide(prs)
    _add_header(slide_bs_lev, "Balance Sheet | Leverage")
    _embed_chart(slide_bs_lev, charts["bs_leverage"])

    # ------------------------------------------------------------------
    # Slide 9 — Cash Flow | OCF vs FCF
    # ------------------------------------------------------------------
    slide_cf_fcf = _add_slide(prs)
    _add_header(slide_cf_fcf, "Cash Flow | OCF vs FCF")
    _embed_chart(slide_cf_fcf, charts["cf_fcf"])

    # ------------------------------------------------------------------
    # Slide 10 — Cash Flow | Trend
    # ------------------------------------------------------------------
    slide_cf_trend = _add_slide(prs)
    _add_header(slide_cf_trend, "Cash Flow | Trend")
    _embed_chart(slide_cf_trend, charts["cf_trend"])

    # ------------------------------------------------------------------
    # Slide 11 — KPI Scorecard
    # ------------------------------------------------------------------
    slide_kpi = _add_slide(prs)
    _add_header(slide_kpi, "KPI Scorecard")
    _embed_chart(slide_kpi, charts["kpi"])

    # ------------------------------------------------------------------
    # Footer pass: iterate all slides now that total count is known
    # ------------------------------------------------------------------
    total = len(prs.slides)
    for idx, sl in enumerate(prs.slides, start=1):
        _add_footer(sl, company, idx, total)

    return prs
